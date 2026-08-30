import subprocess, sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

try:
    from pptx import Presentation
except ImportError:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx', '-q'])
    from pptx import Presentation

prs = Presentation(r'C:\Users\capaw\Downloads\AccentureChallenge\excenture.pptx.pptx')

def extract_text_from_shape(shape, indent=0):
    pad = '  ' * indent
    results = []
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            results.append(f'{pad}TEXT: {repr(text)}')
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    font = run.font
                    size_pt = round(font.size / 12700, 1) if font.size else None
                    try:
                        color = str(font.color.rgb) if font.color and font.color.type else 'inherited'
                    except Exception:
                        color = 'inherited'
                    if run.text.strip():
                        results.append(f'{pad}  run: {repr(run.text)}  bold={font.bold}  sz={size_pt}pt  clr={color}')
    if shape.shape_type == 6:
        try:
            for child in shape.shapes:
                results.extend(extract_text_from_shape(child, indent + 1))
        except Exception:
            pass
    return results

for slide_num, slide in enumerate(prs.slides, 1):
    print(f'\n{"="*70}')
    print(f'SLIDE {slide_num}  |  Layout: {slide.slide_layout.name}')
    print(f'{"="*70}')
    for shape in slide.shapes:
        for line in extract_text_from_shape(shape, 0):
            print(line)
